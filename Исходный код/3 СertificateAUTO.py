#  Импортируем библиотеки
import openpyxl as op # Библиотека для работы с Excel
import pptx
import os # библиотека для файловой работы
import sys # библиотека для работы с системй
from PySide2 import  QtGui, QtWidgets #  библиотеки для создания макета приложеня
from PySide2.QtWidgets import QFileDialog, QMessageBox
from PySide2.QtCore import QCoreApplication
from certificate_form import Ui_MainWindow #  мает приложения
from pathlib import Path
import win32com.client


# Гланвое окно приложения
class MyFirstGuiProgram(QtWidgets.QMainWindow, Ui_MainWindow):
    def __init__(self,
                 parent=None):
        QtWidgets.QMainWindow.__init__(self, parent=parent)
        self.setupUi(self)
        self.btn_cliced()
        self.path_pptx = "Путь к файлу" # Задаем начальные значения переменным
        self.path_excel = "Путь к файлу"
        self.settings_write() # считываем настройки


    def msg_box(self, text_title, text_msg, text_detalic):
        # Функция вывода системных сообщений программы (информирует пользователя об ошибках и событиях в программе)
        error = QMessageBox()
        error.setWindowTitle(text_title)
        error.setText(text_msg)
        error.setIcon(QMessageBox.Warning)
        icon = QtGui.QIcon()
        icon.addPixmap(QtGui.QPixmap("src/ico/certificate.ico"), QtGui.QIcon.Normal, QtGui.QIcon.Off)
        error.setWindowIcon(icon)
        if text_detalic != "":
            error.setDetailedText(text_detalic)
        error.exec_()

    def ppt_to_pdf(self, pptx_path_file,file_name, formatType = 32):

        QCoreApplication.processEvents()  # Обновление формы

        powerpoint = win32com.client.Dispatch("Powerpoint.Application")

        deck = powerpoint.Presentations.Open(pptx_path_file, WithWindow=False)


        try:

            str_path = self.path_save + "/" + f"{file_name}.pdf"
            deck.SaveAs(Path(str_path), formatType)

        except:
            print(27)
        deck.Save()


        deck.Close()

        powerpoint.Quit()


        os.remove(pptx_path_file)



    def create_pptx(self, arr_n1_7):


        # Функция открывает шаблоный файл заполняет данными из таблицы и сохраняет заполненный файл в указанную папку
        QCoreApplication.processEvents()  # Обновление формы
        # Для всех активных тегов записываем их значения из полученного словоря
        if self.chbox_n1.isChecked() == True: n1_text = arr_n1_7["N1"]
        if self.chbox_n2.isChecked() == True: n2_text = arr_n1_7["N2"]
        if self.chbox_n3.isChecked() == True: n3_text = arr_n1_7["N3"]
        if self.chbox_n4.isChecked() == True: n4_text = arr_n1_7["N4"]
        if self.chbox_n5.isChecked() == True: n5_text = arr_n1_7["N5"]
        if self.chbox_n6.isChecked() == True: n6_text = arr_n1_7["N6"]
        if self.chbox_n7.isChecked() == True: n7_text = arr_n1_7["N7"]
        QCoreApplication.processEvents()  # Обновление формы
        prs = pptx.Presentation(self.path_pptx) #открываем шаблон

        QCoreApplication.processEvents()  # Обновление формы

        slide = prs.slides[0] #работаем с первым слайдом
        QCoreApplication.processEvents()  # Обновление формы

        for shape in slide.shapes: # Обходим все  объекты
            QCoreApplication.processEvents()  # Обновление формы

            if not shape.has_text_frame: # Проверяем является ли объект текстовым полем
                continue
            for paragraph in shape.text_frame.paragraphs: # работаем с параграфами текстового поля
                QCoreApplication.processEvents()  # Обновление формы
                for run in paragraph.runs:  # получаем параграфы
                    QCoreApplication.processEvents()  # Обновление формы
                    text_sh = run.text
                    # Проверяем параграф на наличие тегов. Если теги присуствуют, заменяем тег на соотвествующее значение
                    try:
                        if self.chbox_n1.isChecked() == True: text_sh = text_sh.replace("{N1}", str(n1_text))
                    except:
                        pass
                    try:
                        if self.chbox_n2.isChecked() == True: text_sh = text_sh.replace("{N2}", str(n2_text))
                    except:
                        pass
                    try:
                        if self.chbox_n3.isChecked() == True: text_sh = text_sh.replace("{N3}", str(n3_text))
                    except:
                        pass
                    try:
                        if self.chbox_n4.isChecked() == True: text_sh = text_sh.replace("{N4}", str(n4_text))
                    except:
                        pass
                    try:
                        if self.chbox_n5.isChecked() == True: text_sh = text_sh.replace("{N5}", str(n5_text))
                    except:
                        pass
                    try:
                        if self.chbox_n6.isChecked() == True: text_sh = text_sh.replace("{N6}", str(n6_text))
                    except:
                        pass
                    try:
                        if self.chbox_n7.isChecked() == True: text_sh = text_sh.replace("{N7}", str(n7_text))
                    except:
                        pass

                    run.text = text_sh

        QCoreApplication.processEvents()  # Обновление формы
        prs.save(self.path_save + "\\" + f"{n1_text}.pptx") # сохраняем презентацию по указанному пути
        if self.cbox_pdf_2.isChecked() == True: self.ppt_to_pdf(self.path_save + "\\" + f"{n1_text}.pptx", n1_text)

        QCoreApplication.processEvents()  # Обновление формы

    def excel_write(self):
        QCoreApplication.processEvents()  # Обновление формы
        # Функция открывает файл Excel для чтения
        k_start = int(self.ed_start.text()) # Записываем диапазон обработки таблицы
        k_max = (self.ed_end.text()).strip()
        QCoreApplication.processEvents()  # Обновление формы

        self.wb = op.open(filename=self.path_excel, data_only=True) # Открываем таблицу
        self.sheet_read = self.wb[self.cbox_excel_sheets.currentText()] # активируем лист с данными
        # Устанавливаем соотвествие между тегами и колонками таблицы
        n1_coll = self.ed_n1.text()
        n2_coll = self.ed_n2.text()
        n3_coll = self.ed_n3.text()
        n4_coll = self.ed_n4.text()
        n5_coll = self.ed_n5.text()
        n6_coll = self.ed_n6.text()
        n7_coll = self.ed_n7.text()
        QCoreApplication.processEvents()  # Обновление формы

        i=k_start # начало диапазона
        # Получаем данные с файла Excel передаем данные для заполнения в другую функцию
        while (self.sheet_read[f"A{i}"].value) !=None:
            QCoreApplication.processEvents()  # Обновление формы

            arr_n1_7 = {} #  словарь для заполнения данными
            if k_max != "-@" :
                if int(k_max)+1 == i: break # Условие достижения конца таблицы
            # Значение всех активных тегов добавляем в словарь
            if self.chbox_n1.isChecked() == True:
                n1_text = self.sheet_read[f"{n1_coll}{i}"].value
                arr_n1_7["N1"] = n1_text
            if self.chbox_n2.isChecked() == True:
                n2_text = self.sheet_read[f"{n2_coll}{i}"].value
                arr_n1_7["N2"] = n2_text
            if self.chbox_n3.isChecked() == True:
                n3_text = self.sheet_read[f"{n3_coll}{i}"].value
                arr_n1_7["N3"] = n3_text
            if self.chbox_n4.isChecked() == True:
                n4_text = self.sheet_read[f"{n4_coll}{i}"].value
                arr_n1_7["N4"] = n4_text
            if self.chbox_n5.isChecked() == True:
                n5_text = self.sheet_read[f"{n5_coll}{i}"].value
                arr_n1_7["N5"] = n5_text
            if self.chbox_n6.isChecked() == True:
                n6_text = self.sheet_read[f"{n6_coll}{i}"].value
                arr_n1_7["N6"] = n6_text
            QCoreApplication.processEvents()  # Обновление формы
            if self.chbox_n7.isChecked() == True:
                n7_text = self.sheet_read[f"{n7_coll}{i}"].value
                arr_n1_7["N7"] = n7_text
            QCoreApplication.processEvents()  # Обновление формы

            i += 1 # переходим к следующей строке таблицы

            self.create_pptx(arr_n1_7) # переходим к функции заполнения шаблона данными
            QCoreApplication.processEvents()  # Обновление формы


        self.wb.close() # Закрываем таблицу

    def excel_sheets(self):
        QCoreApplication.processEvents()  # Обновление формы
        # Функция считывает названия всех листов книги и записывает их в выпадающий список
        try:
            self.cbox_excel_sheets.clear()
        except:
            pass
        try:
            self.wb = op.open(self.path_excel)
            sheet_name = self.wb.sheetnames
            self.wb.close()
            for i in sheet_name:
                self.cbox_excel_sheets.addItem(i)
        except:
            pass
        QCoreApplication.processEvents()  # Обновление формы


    def open_file_pptx(self):
        QCoreApplication.processEvents()  # Обновление формы
        # Функция открывает диалоговое окно выбора щаблона файла
        self.path_pptx = QFileDialog.getOpenFileName(filter="pptx *pptx")[0]  # Диалоговое окно выбора файла
        QCoreApplication.processEvents()  # Обновление формы
        if self.path_pptx == "":  # Проверка выбран ли файл
            self.path_pptx = "Путь к файлу"
            self.ed_file_pptx.setText(self.path_pptx)
            QCoreApplication.processEvents()  # Обновление формы
            self.msg_box("Ошибка", "Файл не выбран", "")
            return
        else:
            try:
                self.ed_file_pptx.setText(self.path_pptx)
                QCoreApplication.processEvents()  # Обновление формы
            except:
                print("Файл не найден 1 ")
                return
    def open_file_excel(self):
        QCoreApplication.processEvents()  # Обновление формы
        # функция открытия диалогового окна выбора файла Excel с данными
        self.path_excel = QFileDialog.getOpenFileName(filter="EXCEL *xlsx")[0]  # Диалоговое окно выбора файла
        QCoreApplication.processEvents()  # Обновление формы

        if self.path_excel == "":  # Проверка выбран ли файл
            self.path_excel = "Путь к файлу"
            self.ed_file_excel.setText(self.path_excel)
            QCoreApplication.processEvents()  # Обновление формы

            self.msg_box("Ошибка", "Файл не выбран", "")
            return

        else:
            try:
                QCoreApplication.processEvents()  # Обновление формы
                self.ed_file_excel.setText(self.path_excel)
                self.excel_sheets() # вызов функции получения имен листов книги
            except:
                print("Файл не найден 2")
                return

    # Функции активации деактивации полей тегов
    def n1(self):
        if self.ed_n1.isEnabled() == False:
            self.ed_n1.setEnabled(True)
        else:
            self.ed_n1.setEnabled(False)
    def n2(self):
        if self.ed_n2.isEnabled() == False:
            self.ed_n2.setEnabled(True)
        else:
            self.ed_n2.setEnabled(False)
    def n3(self):
        if self.ed_n3.isEnabled() == False:
            self.ed_n3.setEnabled(True)
        else:
            self.ed_n3.setEnabled(False)
    def n4(self):
        if self.ed_n4.isEnabled() == False:
            self.ed_n4.setEnabled(True)
        else:
            self.ed_n4.setEnabled(False)
    def n5(self):
        if self.ed_n5.isEnabled() == False:
            self.ed_n5.setEnabled(True)
        else:
            self.ed_n5.setEnabled(False)
    def n6(self):
        if self.ed_n6.isEnabled() == False:
            self.ed_n6.setEnabled(True)
        else:
            self.ed_n6.setEnabled(False)
    def n7(self):
        if self.ed_n7.isEnabled() == False:
            self.ed_n7.setEnabled(True)
        else:
            self.ed_n7.setEnabled(False)


    def cl_ok(self):
        # Функция вызывает диалоговое окно сохранения файлов
        if self.path_excel == "Путь к файлу":
            self.msg_box("Заполнение сертификатов и дипломов", "Не загружен файл Excel с данными", "")
            return
        if self.path_pptx == "Путь к файлу":
            self.msg_box("Заполнение сертификатов и дипломов", "Не загружен файл  шаблона", "")
            return
        self.path_save = ""
        self.path_save = QFileDialog.getSaveFileName()[0]

        if self.path_save =="":

            self.msg_box("Ошибка", "Не указан путь сохранения файлов", "")
            return
        else:
            try:
                os.mkdir(self.path_save)
                self.excel_write()

            except:
                print("Файл не найден 3")
                return
        self.ed_file_excel.setText("Путь к файлу")
        self.ed_file_pptx.setText("Путь к файлу")
        self.cbox_excel_sheets.clear()
        self.ed_start.setText("2")
        self.ed_end.setText("-@")
        self.msg_box("Заполнение сертификатов и дипломов ", f"Готовые сертификаты и дипломы находятся в папке {self.path_save}", "")

    def setting_save(self):
        # Функция сохраняет настройки соотвествия между тегами и колонками таблицы
        f = open('src/setting/setting.txt', 'w')  # Сохраняем настройки в текстовый файл
        n1 = (self.ed_n1.text()).strip()
        n2 = (self.ed_n2.text()).strip()
        n3 = (self.ed_n3.text()).strip()
        n4 = (self.ed_n4.text()).strip()
        n5 = (self.ed_n5.text()).strip()
        n6 = (self.ed_n6.text()).strip()
        n7 = (self.ed_n7.text()).strip()
        f.write(f"{n1} {n2} {n3} {n4} {n5} {n6} {n7}")
        f.close()
    def settings_write(self):
        # Функция считывания настроек из файла
        f = open('src/setting/setting.txt', 'r')  # открыть файл из рабочей директории в режиме чтения
        for line in f:
            str_text = line.split(" ")
            self.ed_n1.setText(str_text[0])
            self.ed_n2.setText(str_text[1])
            self.ed_n3.setText(str_text[2])
            self.ed_n4.setText(str_text[3])
            self.ed_n5.setText(str_text[4])
            self.ed_n6.setText(str_text[5])
            self.ed_n7.setText(str_text[6])
            break
        f.close()
    def help_b(self):
        # Функция вызова справки
        self.groupBox.setVisible(False)
        self.widget2_2.setVisible(False)
        self.widget2.setVisible(True)
    def shablon_b(self):
        # функция вызова справки по оформлению шаблона
        self.groupBox.setVisible(False)
        self.widget2_2.setVisible(True)
        self.widget2.setVisible(False)
    def menu_b(self):
        # функция закрытия справок
        self.groupBox.setVisible(True)
        self.widget2_2.setVisible(False)
        self.widget2.setVisible(False)
    def btn_cliced(self):
        # функция отслеживания на нажатия кнопок и вызова соотвествующих функция
        self.btn_help.clicked.connect(lambda: self.help_b())
        self.btn_shablon.clicked.connect(lambda: self.shablon_b())
        self.btn_open_file_pptx.clicked.connect(lambda: self.open_file_pptx())
        self.btn_open_file_excel.clicked.connect(lambda: self.open_file_excel())
        self.btn_seting_save.clicked.connect(lambda: self.setting_save())
        self.btn_ok.clicked.connect(lambda:self.cl_ok() )
        self.chbox_n1.clicked.connect(lambda: self.n1())
        self.chbox_n2.clicked.connect(lambda: self.n2())
        self.chbox_n3.clicked.connect(lambda: self.n3())
        self.chbox_n4.clicked.connect(lambda: self.n4())
        self.chbox_n5.clicked.connect(lambda: self.n5())
        self.chbox_n6.clicked.connect(lambda: self.n6())
        self.chbox_n7.clicked.connect(lambda: self.n7())
        self.btn_close_help.clicked.connect(lambda:self.menu_b())
        self.btn_close_sh.clicked.connect(lambda:self.menu_b())


if __name__ == '__main__':
    # Запуск программы начальные настройки
    app = QtWidgets.QApplication(sys.argv)
    app.setStyle("Fusion")
    MainWindow = MyFirstGuiProgram()
    MainWindow.show()
    sys.exit(app.exec_())



